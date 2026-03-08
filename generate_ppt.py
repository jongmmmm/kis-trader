"""kis-trader 코드 분석 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0x7C, 0x4D, 0xFF)
ACCENT3 = RGBColor(0x00, 0xE6, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
RED = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW = RGBColor(0xFF, 0xE6, 0x6D)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=BG_MID, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)

        # bullet marker
        run_b = p.add_run()
        run_b.text = "  >  "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.bold = True

        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tf


def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT, bg_color=BG_MID):
    add_shape_bg(slide, left, top, width, height, bg_color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.45),
                title, font_size=16, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                    body_items, font_size=13, color=LIGHT_GRAY, bullet_color=title_color)


# ════════════════════════════════════════════
# SLIDE 1: 타이틀
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 가운데 정렬 제목
add_textbox(slide, Inches(0), Inches(1.5), Inches(13.333), Inches(1),
            "KIS-TRADER", font_size=54, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(2.6), Inches(13.333), Inches(0.8),
            "한국투자증권 Open API 기반 자동매매 시스템", font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.6),
            "코드 분석 보고서", font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 하단 태그들
tags = "Flask  |  SQLAlchemy  |  APScheduler  |  REST API  |  SSE  |  ML (scikit-learn)  |  RAG AI Analyzer"
add_textbox(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.5),
            tags, font_size=14, color=ACCENT2, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 2: 목차
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "목차 (Table of Contents)", font_size=32, color=ACCENT, bold=True)

toc = [
    "1. 프로젝트 개요 & 기술 스택",
    "2. 전체 아키텍처 & 파일 구조",
    "3. 핵심 모듈: KIS API 래퍼 (kis_api.py)",
    "4. 데이터 모델 (models.py)",
    "5. 매매 전략 엔진 (strategies/)",
    "6. RAG AI 분석기 (ai_analyzer.py)",
    "7. 스케줄러 & 자동화 (scheduler.py)",
    "8. 웹 API & 라우터 (routers/)",
    "9. 실시간 경매 알림 (SSE)",
    "10. 시스템 흐름도 (End-to-End Flow)",
    "11. 주요 특징 & 개선 포인트",
]
add_bullet_list(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(5.5),
                toc, font_size=20, color=WHITE, bullet_color=ACCENT2)

# ════════════════════════════════════════════
# SLIDE 3: 프로젝트 개요 & 기술 스택
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "1. 프로젝트 개요 & 기술 스택", font_size=32, color=ACCENT, bold=True)

add_textbox(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.8),
            "한국투자증권(KIS) Open API를 활용한 주식 자동매매 웹 애플리케이션.\n"
            "모의투자/실전투자 모드 전환, 4가지 매매 전략, AI 분석, 실시간 경매 알림을 제공합니다.",
            font_size=17, color=LIGHT_GRAY)

# 기술 스택 카드들
cards = [
    ("Backend", ["Flask 3.1 (웹 프레임워크)", "Flask-SQLAlchemy (ORM)", "SQLite (데이터베이스)", "APScheduler (스케줄링)"], ACCENT),
    ("Data & ML", ["pandas (데이터 처리)", "ta (기술적 지표 라이브러리)", "scikit-learn (RandomForest)", "joblib (모델 직렬화)"], ACCENT3),
    ("API & 통신", ["requests (KIS API 호출)", "SSE (Server-Sent Events)", "REST API (JSON)", "python-dotenv (환경변수)"], ACCENT2),
    ("운영 모드", ["모의투자 (paper) 모드", "실전투자 (real) 모드", "런타임 모드 전환 지원", "별도 API 키/계좌 분리"], ORANGE),
]

for i, (title, items, col) in enumerate(cards):
    left = Inches(0.4 + i * 3.2)
    add_card(slide, left, Inches(2.5), Inches(3.0), Inches(4.2), title, items, title_color=col)

# ════════════════════════════════════════════
# SLIDE 4: 전체 아키텍처 & 파일 구조
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "2. 전체 아키텍처 & 파일 구조", font_size=32, color=ACCENT, bold=True)

# 파일 트리
tree_items = [
    "run.py               - 애플리케이션 진입점 (port 6001)",
    "app.py               - Flask 앱 팩토리 (create_app)",
    "config.py            - 환경 설정 (API 키, URL, 스케줄러 간격)",
    "db.py                - SQLAlchemy 인스턴스",
    "models.py            - DB 모델 5개 (Strategy, Order, AuctionAlert ...)",
    "kis_api.py           - KIS Open API 래퍼 (7개 함수)",
    "scheduler.py         - APScheduler 초기화 & 4개 Job 등록",
    "routers/             - Blueprint 5개 (dashboard, strategies, orders, auction, settings)",
    "strategies/          - 매매 전략 모듈 (base, ma, rsi_macd, condition, ml, ai_analyzer, runner)",
]
add_bullet_list(slide, Inches(0.5), Inches(1.2), Inches(7), Inches(5.5),
                tree_items, font_size=15, color=WHITE, bullet_color=ACCENT)

# 아키텍처 다이어그램 (텍스트 기반)
add_shape_bg(slide, Inches(8), Inches(1.2), Inches(5), Inches(5.5), BG_MID)
add_textbox(slide, Inches(8.2), Inches(1.3), Inches(4.5), Inches(0.4),
            "Architecture Flow", font_size=16, color=ACCENT, bold=True)

arch_lines = [
    "[Browser/Client]",
    "       |",
    "  Flask Web Server (run.py:6001)",
    "    |-- /api/balance",
    "    |-- /api/market-prices",
    "    |-- /api/strategies (CRUD)",
    "    |-- /api/auction/stream (SSE)",
    "    |-- /api/settings/mode",
    "       |",
    "  APScheduler (Background)",
    "    |-- run_strategies (60s)",
    "    |-- check_auction (30s)",
    "    |-- expire_alerts (60s)",
    "    |-- retrain_ml (03:00 daily)",
    "       |",
    "  KIS Open API (REST)",
    "    |-- 현재가/일봉/주문/잔고",
]
add_textbox(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(4.8),
            "\n".join(arch_lines), font_size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 5: KIS API 래퍼
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "3. 핵심 모듈: KIS API 래퍼 (kis_api.py)", font_size=32, color=ACCENT, bold=True)

add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
            "한국투자증권 Open API를 추상화한 Python 래퍼. 토큰 자동 갱신, 모의/실전 URL 자동 전환.",
            font_size=16, color=LIGHT_GRAY)

api_funcs = [
    ("get_token(mode)", "OAuth 토큰 발급/갱신. DB 캐시로 만료 5분 전 자동 재발급", ACCENT),
    ("get_current_price(code, mode)", "종목 현재가 조회 (가격/등락률/거래량/고가/저가/시가)", ACCENT3),
    ("get_daily_ohlcv(code, mode, count, period)", "일봉/주봉/월봉 OHLCV 데이터. 최대 100개", ACCENT2),
    ("place_order(code, type, price, qty, mode)", "매수/매도 주문 실행. 지정가(00)/시장가(01) 지원", ORANGE),
    ("get_index_price(index_code, mode)", "업종 지수 조회 (코스피 0001 / 코스닥 1001)", ACCENT),
    ("get_volume_rank(mode)", "거래량 상위 30종목 조회", ACCENT3),
    ("get_balance(mode)", "계좌 잔고 조회 (보유종목/평단/수익률)", ACCENT2),
]

for i, (name, desc, col) in enumerate(api_funcs):
    y = Inches(1.7 + i * 0.75)
    add_shape_bg(slide, Inches(0.5), y, Inches(12.3), Inches(0.65), BG_MID)
    add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(4.5), Inches(0.3),
                name, font_size=14, color=col, bold=True)
    add_textbox(slide, Inches(5.3), y + Inches(0.05), Inches(7.3), Inches(0.55),
                desc, font_size=13, color=LIGHT_GRAY)

# 핵심 특징
add_shape_bg(slide, Inches(0.5), Inches(7.0) - Inches(0.8), Inches(12.3), Inches(0.6), BG_MID)
add_textbox(slide, Inches(0.7), Inches(7.0) - Inches(0.75), Inches(12), Inches(0.5),
            "핵심: _credentials(mode)로 paper/real 자동 분기  |  _headers()에 tr_id 매핑  |  KisToken 모델로 DB 캐싱",
            font_size=13, color=YELLOW, bold=True)

# ════════════════════════════════════════════
# SLIDE 6: 데이터 모델
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "4. 데이터 모델 (models.py)", font_size=32, color=ACCENT, bold=True)

models = [
    ("Strategy", [
        "id, name, stock_code, stock_name",
        "strategy_type: ma | rsi_macd | condition | ml",
        "params: JSON (전략별 파라미터)",
        "is_active, mode (paper/real)"
    ], ACCENT),
    ("Order", [
        "id, strategy_id (FK), stock_code",
        "order_type: buy | sell",
        "status: pending | submitted | filled | cancelled",
        "trigger: auto | auction_manual",
        "kis_order_no: KIS 주문번호"
    ], ACCENT3),
    ("AuctionAlert", [
        "strategy_id, stock_code, stock_name",
        "suggested_action/price/qty",
        "ai_confidence (0~100), ai_score (-100~+100)",
        "ai_summary, ai_factors (JSON)",
        "user_decision, expires_at"
    ], ORANGE),
    ("PortfolioSnapshot", [
        "mode, stock_code, quantity",
        "avg_price, current_price",
        "snapped_at (스냅샷 시점)"
    ], ACCENT2),
    ("KisToken", [
        "mode (paper/real, unique)",
        "access_token, expires_at",
        "OAuth 토큰 DB 캐싱용"
    ], RED),
]

for i, (name, fields, col) in enumerate(models):
    col_idx = i % 3
    row_idx = i // 3
    left = Inches(0.3 + col_idx * 4.3)
    top = Inches(1.2 + row_idx * 3.3)
    h = Inches(3.0) if row_idx == 0 else Inches(2.5)
    add_card(slide, left, top, Inches(4.1), h, name, fields, title_color=col)

# ════════════════════════════════════════════
# SLIDE 7: 매매 전략 엔진
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "5. 매매 전략 엔진 (strategies/)", font_size=32, color=ACCENT, bold=True)

add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
            "BaseStrategy 추상 클래스를 상속한 4가지 전략. should_buy() / should_sell() 인터페이스로 통일.",
            font_size=16, color=LIGHT_GRAY)

strats = [
    ("MAStrategy (이동평균)", [
        "단기/장기 이동평균선 계산 (기본 5일/20일)",
        "골든크로스 -> 매수 신호",
        "데드크로스 -> 매도 신호",
        "pandas rolling().mean() 활용"
    ], ACCENT, "ma_strategy.py"),
    ("RsiMacdStrategy (RSI+MACD)", [
        "RSI 과매도(30) + MACD 상향돌파 -> 매수",
        "RSI 과매수(70) + MACD 하향돌파 -> 매도",
        "ta 라이브러리로 지표 계산",
        "파라미터 커스터마이징 가능"
    ], ACCENT3, "rsi_macd.py"),
    ("ConditionStrategy (조건식)", [
        "사용자 정의 조건식 평가",
        "indicator + operator + value 구조",
        ">, <, >=, <=, == 연산자 지원",
        "action: buy / sell / both 설정"
    ], ACCENT2, "condition.py"),
    ("MLStrategy (머신러닝)", [
        "RandomForest 분류기 (scikit-learn)",
        "피처: MA5/20/60, RSI, MACD, 거래량비율, 수익률",
        "매일 03시 자동 재학습 (CronTrigger)",
        "확률 임계값으로 매수/매도 판단 (0.65/0.35)"
    ], ORANGE, "ml_strategy.py"),
]

for i, (name, items, col, filename) in enumerate(strats):
    left = Inches(0.3 + (i % 2) * 6.5)
    top = Inches(1.6 + (i // 2) * 2.8)
    add_card(slide, left, top, Inches(6.3), Inches(2.6), f"{name}  [{filename}]", items, title_color=col)

# BaseStrategy 공통 기능
add_shape_bg(slide, Inches(0.3), Inches(7.0) - Inches(0.7), Inches(12.7), Inches(0.55), BG_MID)
add_textbox(slide, Inches(0.5), Inches(7.0) - Inches(0.65), Inches(12.3), Inches(0.5),
            "BaseStrategy 공통: get_quantity() | check_stop_loss(stop_loss_pct, 기본 -5%) | check_take_profit(take_profit_pct, 기본 +10%)",
            font_size=14, color=YELLOW, bold=True)

# ════════════════════════════════════════════
# SLIDE 8: RAG AI 분석기
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "6. RAG AI 분석기 (ai_analyzer.py)", font_size=32, color=ACCENT, bold=True)

add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
            "8가지 팩터를 분석하고 가중 합산하여 매수/매도/관망을 추천. 신뢰도와 근거 텍스트를 자동 생성.",
            font_size=16, color=LIGHT_GRAY)

factors = [
    ("이동평균 추세", "20%", "골든/데드크로스, 정배열/역배열, 60일선 위치", ACCENT),
    ("RSI", "15%", "과매수/과매도 (20/30/70/80), RSI 추세", ACCENT3),
    ("MACD", "15%", "히스토그램 양/음 전환, 매수/매도세 확대", ACCENT2),
    ("거래량", "15%", "거래량 폭발(3x), 증가(1.5x), 부진(0.5x)", ORANGE),
    ("모멘텀", "10%", "5일/20일 수익률, 과열 경고, 저점 반등", RED),
    ("변동성", "10%", "ATR 기반, 고변동(>5%), 안정(<1.5%)", YELLOW),
    ("캔들 패턴", "8%", "망치형, 교수형, 도지, 3일 연속 양/음봉", ACCENT),
    ("당일 등락", "7%", "급등/급락 분석, 장중 고/저점 위치", ACCENT3),
]

for i, (name, weight, desc, col) in enumerate(factors):
    y = Inches(1.6 + i * 0.62)
    add_shape_bg(slide, Inches(0.5), y, Inches(8.5), Inches(0.55), BG_MID)
    add_textbox(slide, Inches(0.7), y + Inches(0.05), Inches(2.0), Inches(0.3),
                name, font_size=14, color=col, bold=True)
    add_textbox(slide, Inches(2.8), y + Inches(0.05), Inches(0.8), Inches(0.3),
                weight, font_size=14, color=YELLOW, bold=True)
    add_textbox(slide, Inches(3.7), y + Inches(0.05), Inches(5.2), Inches(0.45),
                desc, font_size=12, color=LIGHT_GRAY)

# 판단 로직 카드
add_card(slide, Inches(9.3), Inches(1.6), Inches(3.8), Inches(5.0),
         "판단 로직",
         [
             "각 팩터 점수: -100 ~ +100",
             "가중 합산 -> total_score",
             "score >= 20 -> 매수 (buy)",
             "score <= -20 -> 매도 (sell)",
             "-20 < score < 20 -> 관망 (hold)",
             "신뢰도 = |score|*0.7 + 일치도*30",
             "최대 95%까지 캡",
             "요약 텍스트 자동 생성 (5줄)",
         ], title_color=ACCENT2)

# ════════════════════════════════════════════
# SLIDE 9: 스케줄러 & 자동화
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "7. 스케줄러 & 자동화 (scheduler.py)", font_size=32, color=ACCENT, bold=True)

add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
            "APScheduler BackgroundScheduler로 4개의 자동 Job을 관리. Asia/Seoul 타임존.",
            font_size=16, color=LIGHT_GRAY)

jobs = [
    ("run_strategies", "60초 간격 (IntervalTrigger)", [
        "장중(09:00~15:30) & 비경매 시간에만 실행",
        "활성 전략 순회 -> 지표 계산 -> 매매 신호 판단",
        "보유 종목 손절/익절 우선 체크",
        "주문 실행 후 Order 레코드 기록"
    ], ACCENT),
    ("check_auction_and_alert", "30초 간격 (IntervalTrigger)", [
        "경매 시간(08:00~09:00, 15:20~15:30)에만 실행",
        "활성 전략별 AI 분석 수행",
        "AuctionAlert 생성 -> SSE로 브로드캐스트",
        "사용자 결정 대기 (buy/sell/pass)"
    ], ORANGE),
    ("expire_undecided_alerts", "60초 간격 (IntervalTrigger)", [
        "만료 시간 초과 미결정 알림 자동 pass 처리",
        "DB 일괄 업데이트"
    ], ACCENT2),
    ("retrain_ml_models", "매일 03:00 (CronTrigger)", [
        "ML 전략 종목 120일 데이터로 재학습",
        "RandomForest 모델 joblib 저장",
        "피처: MA, RSI, MACD, 거래량비율, 수익률"
    ], ACCENT3),
]

for i, (name, trigger, items, col) in enumerate(jobs):
    left = Inches(0.3 + (i % 2) * 6.5)
    top = Inches(1.6 + (i // 2) * 2.8)
    add_card(slide, left, top, Inches(6.3), Inches(2.6),
             f"{name}  ({trigger})", items, title_color=col)

# ════════════════════════════════════════════
# SLIDE 10: 웹 API & 라우터
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "8. 웹 API & 라우터 (routers/)", font_size=32, color=ACCENT, bold=True)

routers = [
    ("Dashboard (dashboard.py)", [
        "GET /                    -> 대시보드 HTML",
        "GET /api/balance         -> 계좌 잔고 조회",
        "GET /api/market-prices   -> 코스피/코스닥 + 종목 시세",
        "GET /api/market-ranking  -> 거래량 상위 종목",
        "GET /api/stock-price/<code> -> 개별 종목 시세",
    ], ACCENT),
    ("Strategies (strategies.py)", [
        "GET  /strategies/                      -> 전략 관리 HTML",
        "GET  /api/strategies                    -> 전략 목록",
        "POST /api/strategies                    -> 전략 등록",
        "PUT  /api/strategies/<id>               -> 전략 수정",
        "DELETE /api/strategies/<id>             -> 전략 삭제",
        "POST /api/strategies/<id>/toggle        -> 활성/비활성 토글",
        "GET  /api/strategies/<id>/chart          -> 차트 데이터 (OHLCV + 지표)",
        "GET  /api/strategies/<id>/ai-analyze     -> AI 분석 수동 실행",
    ], ACCENT3),
    ("Orders (orders.py)", [
        "GET /orders/         -> 주문 내역 HTML",
        "GET /api/orders      -> 주문 목록 (mode/trigger 필터)",
    ], ACCENT2),
    ("Settings (settings.py)", [
        "GET  /api/settings/mode         -> 현재 모드 조회",
        "POST /api/settings/mode         -> paper/real 전환",
        "POST /api/settings/token-refresh -> 토큰 캐시 초기화",
    ], ORANGE),
]

for i, (name, items, col) in enumerate(routers):
    if i < 2:
        left = Inches(0.3 + i * 6.5)
        top = Inches(1.1)
        h = Inches(3.2) if i == 0 else Inches(4.0)
    else:
        left = Inches(0.3 + (i - 2) * 6.5)
        top = Inches(5.3) if i == 2 else Inches(5.3)
        h = Inches(1.8)
    w = Inches(6.3)
    add_card(slide, left, top, w, h, name, items, title_color=col)

# ════════════════════════════════════════════
# SLIDE 11: 실시간 경매 알림 (SSE)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "9. 실시간 경매 알림 (SSE) - auction.py", font_size=32, color=ACCENT, bold=True)

add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
            "Server-Sent Events를 활용한 실시간 경매 알림 시스템. 반자동 의사결정 지원.",
            font_size=16, color=LIGHT_GRAY)

# SSE 흐름
add_card(slide, Inches(0.3), Inches(1.7), Inches(6.3), Inches(5.0),
         "SSE 스트림 구조 (auction.py)",
         [
             "_sse_clients: List[Queue] - 연결된 클라이언트 큐 관리",
             "threading.Lock으로 동시성 보호",
             "GET /api/auction/stream -> SSE 스트림 연결",
             "broadcast_auction(alert_id) -> 모든 클라이언트에 푸시",
             "Queue(maxsize=20) per client",
             "30초 간격 ping (: ping) 으로 연결 유지",
             "GeneratorExit 시 자동 클린업",
             "Dead queue 감지 후 제거",
         ], title_color=ACCENT)

# 의사결정 흐름
add_card(slide, Inches(6.8), Inches(1.7), Inches(6.2), Inches(5.0),
         "반자동 의사결정 흐름",
         [
             "1. 스케줄러 -> check_auction_and_alert() 호출",
             "2. 활성 전략별 AI 분석 (analyze_stock)",
             "3. AuctionAlert 생성 (DB 저장)",
             "4. broadcast_auction() -> SSE 푸시",
             "5. 브라우저에서 알림 수신",
             "6. 사용자가 buy/sell/pass 결정",
             "7. POST /api/auction/decide/<id>",
             "8. buy/sell 시 KIS API 주문 실행",
             "9. Order 레코드 자동 생성 (trigger=auction_manual)",
             "10. 미결정 알림은 자동 만료 (pass)"
         ], title_color=ORANGE)

# ════════════════════════════════════════════
# SLIDE 12: 시스템 흐름도
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "10. 시스템 흐름도 (End-to-End Flow)", font_size=32, color=ACCENT, bold=True)

# 자동매매 흐름
add_card(slide, Inches(0.3), Inches(1.2), Inches(4.1), Inches(5.5),
         "자동매매 흐름 (장중 09:00~15:30)",
         [
             "1. APScheduler 60초 마다 트리거",
             "2. is_market_hours() 체크",
             "3. DB에서 활성 전략 조회",
             "4. KIS API로 OHLCV + 현재가 조회",
             "5. Strategy 엔진 실행",
             "   - 보유 종목: 손절/익절 체크",
             "   - should_sell() 체크",
             "   - should_buy() 체크",
             "6. place_order() 실행",
             "7. Order 레코드 DB 기록",
         ], title_color=ACCENT)

# 경매 알림 흐름
add_card(slide, Inches(4.6), Inches(1.2), Inches(4.1), Inches(5.5),
         "경매 알림 흐름 (08:00~09:00 / 15:20~15:30)",
         [
             "1. APScheduler 30초 마다 트리거",
             "2. is_auction_time() 체크",
             "3. 활성 전략별 AI 분석 수행",
             "4. 8가지 팩터 가중 합산",
             "5. AuctionAlert 생성",
             "6. SSE로 클라이언트 브로드캐스트",
             "7. 사용자 결정 대기",
             "8. 결정 시 주문 실행 or pass",
             "9. 미결정 시 자동 만료",
         ], title_color=ORANGE)

# ML 학습 흐름
add_card(slide, Inches(8.9), Inches(1.2), Inches(4.1), Inches(5.5),
         "ML 모델 학습 흐름 (매일 03:00)",
         [
             "1. CronTrigger(hour=3) 실행",
             "2. ML 전략 종목 조회",
             "3. KIS API로 120일 OHLCV",
             "4. 피처 엔지니어링",
             "   - MA5/20/60, RSI, MACD diff",
             "   - 거래량 비율, 1일/5일 수익률",
             "5. 타겟: 다음날 상승 여부 (0/1)",
             "6. RandomForest(100 trees) 학습",
             "7. joblib으로 모델 저장",
             "8. 다음 장중에 예측 사용",
         ], title_color=ACCENT3)

# ════════════════════════════════════════════
# SLIDE 13: 주요 특징 & 개선 포인트
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            "11. 주요 특징 & 개선 포인트", font_size=32, color=ACCENT, bold=True)

add_card(slide, Inches(0.3), Inches(1.2), Inches(6.3), Inches(3.0),
         "주요 특징",
         [
             "모의/실전 투자 모드 런타임 전환",
             "4가지 매매 전략 + AI 분석기 (8팩터)",
             "SSE 기반 실시간 경매 알림 (반자동 매매)",
             "ML 모델 자동 재학습 (새벽 3시)",
             "손절(-5%) / 익절(+10%) 자동 관리",
             "Flask Blueprint로 모듈 분리",
             "토큰 DB 캐싱으로 API 호출 최적화",
         ], title_color=ACCENT3)

add_card(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(3.0),
         "개선 포인트",
         [
             "SQLite -> PostgreSQL 전환 (동시성)",
             "비동기 처리 (asyncio/Celery)",
             "웹소켓 실시간 시세 (KIS WebSocket API)",
             "백테스팅 모듈 추가",
             "사용자 인증/권한 관리",
             "Docker 컨테이너화 & CI/CD",
             "모니터링/로깅 강화 (Prometheus/Grafana)",
         ], title_color=RED)

# 코드 통계
add_card(slide, Inches(0.3), Inches(4.5), Inches(4.0), Inches(2.5),
         "코드 통계",
         [
             "Python 소스 파일: 16개",
             "DB 모델: 5개 테이블",
             "API 엔드포인트: 15개",
             "매매 전략: 4종류",
             "AI 분석 팩터: 8개",
             "스케줄러 Job: 4개",
         ], title_color=ACCENT)

add_card(slide, Inches(4.5), Inches(4.5), Inches(4.1), Inches(2.5),
         "디자인 패턴",
         [
             "Factory Pattern (create_app)",
             "Strategy Pattern (BaseStrategy 상속)",
             "Observer Pattern (SSE broadcast)",
             "Template Method (should_buy/sell)",
             "Caching (토큰 DB 캐싱)",
         ], title_color=ACCENT2)

add_card(slide, Inches(8.8), Inches(4.5), Inches(4.2), Inches(2.5),
         "핵심 의존성",
         [
             "Flask 3.1 + SQLAlchemy",
             "APScheduler 3.10",
             "pandas + ta (기술적 지표)",
             "scikit-learn (RandomForest)",
             "requests (KIS API 통신)",
         ], title_color=ORANGE)

# ════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════
output_path = "/home/pois/kis-trader/kis-trader-analysis.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
